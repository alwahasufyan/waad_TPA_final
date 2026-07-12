#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
=============================================================================
  ingest.py  |  محرك ابتلاع قوائم الأسعار متعدّد الصيغ
=============================================================================
يقرأ أي قائمة أسعار واردة (Excel / CSV / PDF / PowerPoint) ويحوّلها إلى جدول
موحّد جاهز للمطابقة، بصرف النظر عن فوضى الترتيب أو غياب العناوين.

الواجهة الوحيدة المطلوبة من السكربت الرئيسي:

    load_records(path) -> pandas.DataFrame
        الأعمدة: name | name_alt | code | price | source

كيف يكتشف الأعمدة دون عناوين؟
  يصنّف كل عمود حسب خصائصه (نسبة الأرقام، نمط الأكواد، طول النص، تسلسل ترقيم
  الصفوف) فيحدّد: عمود الاسم، عمود الاسم البديل، عمود السعر، عمود الكود،
  ويتجاهل عمود ترقيم الصفوف وعناوين الأقسام والصفوف/الأعمدة الفارغة.

PDF:
  يستخرج الجداول من كل الصفحات، وإن لم توجد جداول واضحة يحلّل الأسطر نصّيًا
  (الاسم + السعر في آخر السطر). ويصحّح اتجاه النص العربي المعكوس تلقائيًا.

PowerPoint:
  يقرأ الجداول؛ وإن كانت البيانات في مربّعات نصّية منفصلة يزاوج كل سعر بأقرب
  اسم حسب موضعه على الشريحة (تقريبي — يُراجَع بشريًا).
=============================================================================
"""
import os
import re

import pandas as pd

OUT_COLS = ["name", "name_alt", "code", "price", "source"]

# ---------------------------------------------------------------------------
# أدوات نصّية مساعدة
# ---------------------------------------------------------------------------
_AR = re.compile(r"[\u0600-\u06FF]")
_LATIN = re.compile(r"[A-Za-z]")
_CODE_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9\-_/.]{0,18}$")
_INT_RE = re.compile(r"^\d+$")


def _clean(v) -> str:
    """يحوّل أي قيمة خلية إلى نص نظيف؛ الفراغ/nan/None -> ''."""
    if v is None:
        return ""
    s = str(v).replace("\u00a0", " ").strip()
    return "" if s.lower() in ("nan", "none", "nat") else s


def has_arabic(s: str) -> bool:
    return bool(_AR.search(s))


def has_latin(s: str) -> bool:
    return bool(_LATIN.search(s))


def has_letter(s: str) -> bool:
    return has_arabic(s) or has_latin(s)


def is_price(s: str):
    """يعيد السعر كنص رقمي إن كانت الخلية رقمًا (مع فواصل آلاف)، وإلا None."""
    s = s.strip().replace("٫", ".").replace("،", ",")
    if not s:
        return None
    s2 = re.sub(r"[,\s]", "", s)
    s2 = s2.replace("%", "")
    if not s2:
        return None
    try:
        f = float(s2)
    except ValueError:
        return None
    if f < 0:
        return None
    # رقم صحيح بلا كسور يُعرض بلا منزلة عشرية
    return str(int(f)) if f == int(f) else str(f)


def is_pure_int(s: str) -> bool:
    return bool(_INT_RE.match(s.strip()))


def is_code(s: str) -> bool:
    """كود مستشفى: رمز بلا مسافات يحوي حروفًا لاتينية (مثل EW-001, RETIC, EM1)."""
    s = s.strip()
    if not s or " " in s:
        return False
    if not _CODE_RE.match(s):
        return False
    if s.isdigit():                      # رقم صرف -> سعر/ترقيم لا كود
        return False
    has_d = any(c.isdigit() for c in s)
    has_a = any(c.isascii() and c.isalpha() for c in s)
    has_sep = any(c in "-_/." for c in s)
    if not has_a:                        # لا حروف لاتينية -> ليس كود مستشفى نمطي
        return False
    return has_d or has_sep or len(s) <= 6


# ---------------------------------------------------------------------------
# تصحيح اتجاه النص العربي المعكوس (شائع في PDF)
# ---------------------------------------------------------------------------
def _ar_orientation_score(s: str) -> int:
    """درجة «صحّة اتجاه» النص العربي: تعتمد على قواعد موضع الحروف.

    - التاء المربوطة (ة) و(ه) تأتيان عادةً في آخر الكلمة (+).
    - أداة التعريف (ال) تأتي في بداية الكلمة (+).
    - بداية كلمة بـ(ة) دليل قوي على الانعكاس (-).
    - همزة القطع/الألف في البداية شائعة (+)، والألف/اللام في النهاية أقل (-).
    """
    sc = 0
    for w in s.split():
        if not w:
            continue
        if w[-1] in "ةه":
            sc += 2
        if w.startswith("ال"):
            sc += 2
        if w[0] in "ةه":
            sc -= 2
        if w[-1] in "ال":
            sc -= 1
        if w[0] in "اأإآ":
            sc += 1
    return sc


def fix_arabic(s: str) -> str:
    """يكشف انعكاس اتجاه النص العربي (الشائع في PDF) ويصحّحه فقط عند اللزوم."""
    if not s or not has_arabic(s):
        return s
    rev = s[::-1]
    if _ar_orientation_score(rev) > _ar_orientation_score(s):
        return rev
    return s


# ===========================================================================
# مُصنِّف الأعمدة — قلب المحرك
# ===========================================================================
_HEADER_KEYS = {
    # "ignore": أعمدة وصفية لا تدخل في المطابقة (مصدر الملف، ملاحظات، تسلسل)
    "ignore": ["المصدر", "مصدر", "source", "ملاحظات", "ملاحظه", "ملاحظة",
               "notes", "note", "remark", "التسلسل", "تسلسل"],
    "code": ["code", "كود", "الكود", "رمز", "no.", "no", "م", "رقم"],
    "price": ["price", "سعر", "السعر", "التكلفة", "المبلغ", "قيمة", "القيمة",
              "تكلفة", "amount", "cost", "fee", "rate"],
    "name": ["name", "الاسم", "اسم", "البيان", "بيان", "procedure", "service",
             "الخدمة", "الخدمه", "description", "desc", "test", "التحليل",
             "الفحص", "العملية", "البند", "تفاصيل"],
}


def _match_header_role(cell: str):
    c = cell.strip().lower()
    if not c:
        return None
    for role, keys in _HEADER_KEYS.items():
        for k in keys:
            if k == c or (len(k) >= 3 and k in c):
                return role
    return None


def _detect_header_row(grid, ncols):
    """يبحث عن صف عناوين ضمن أول 15 صفًا (خليتان فأكثر تطابقان مفاتيح معروفة).

    يعيد (رقم الصف, خريطة الأدوار, مجموعة الأعمدة المتجاهَلة).
    """
    best_row, best_map, best_ign, best_hits = None, None, set(), 1
    for r in range(min(15, len(grid))):
        roles, ign = {}, set()
        for c in range(ncols):
            cell = grid[r][c] if c < len(grid[r]) else ""
            role = _match_header_role(cell)
            if role == "ignore":
                ign.add(c)
            elif role and role not in roles:
                roles[role] = c
        if len(roles) + len(ign) > best_hits:
            best_hits = len(roles) + len(ign)
            best_row, best_map, best_ign = r, roles, ign
    return best_row, (best_map or {}), best_ign


def _col_cells(grid, c, start=0):
    return [(_clean(grid[r][c]) if c < len(grid[r]) else "")
            for r in range(start, len(grid))]


def _is_rownum_col(cells) -> bool:
    vals = []
    for c in cells:
        c = c.strip()
        if c == "":
            continue
        if is_pure_int(c):
            vals.append(int(c))
        else:
            return False
    if len(vals) < 3:
        return False
    inc = sum(1 for i in range(1, len(vals)) if vals[i] == vals[i - 1] + 1)
    return inc >= 0.6 * (len(vals) - 1)


def _name_score(cells) -> float:
    """مقياس «كم هذا العمود اسمي»: نسبة الخلايا الحاملة لحروف × متوسط الطول."""
    nonempty = [c for c in cells if c]
    if not nonempty:
        return 0.0
    lettered = [c for c in nonempty if has_letter(c)]
    if not lettered:
        return 0.0
    frac = len(lettered) / len(nonempty)
    avg_len = sum(len(c) for c in lettered) / len(lettered)
    spacey = sum(1 for c in lettered if " " in c) / len(lettered)
    return frac * avg_len * (1 + spacey)


def _price_ratio(cells) -> float:
    nonempty = [c for c in cells if c]
    if not nonempty:
        return 0.0
    return sum(1 for c in nonempty if is_price(c) is not None) / len(nonempty)


def _code_ratio(cells) -> float:
    nonempty = [c for c in cells if c]
    if not nonempty:
        return 0.0
    return sum(1 for c in nonempty if is_code(c)) / len(nonempty)


def extract_records_from_grid(grid, source):
    """يحوّل شبكة خلايا (صفوف × أعمدة) إلى سجلات name/name_alt/code/price."""
    if not grid:
        return []
    # 1) نظّف وأزل الصفوف والأعمدة الفارغة تمامًا
    grid = [[_clean(x) for x in row] for row in grid]
    ncols = max((len(r) for r in grid), default=0)
    grid = [r + [""] * (ncols - len(r)) for r in grid]
    keep_cols = [c for c in range(ncols)
                 if any(grid[r][c] for r in range(len(grid)))]
    if not keep_cols:
        return []
    grid = [[row[c] for c in keep_cols] for row in grid]
    grid = [row for row in grid if any(row)]
    ncols = len(keep_cols)
    if not grid:
        return []

    # 2) اكتشف صف العناوين (إن وُجد) لتحديد بداية البيانات وأدوار الأعمدة
    hdr_row, hdr_map, ignore_cols = _detect_header_row(grid, ncols)
    start = (hdr_row + 1) if hdr_row is not None else 0

    cols = {c: _col_cells(grid, c, start) for c in range(ncols)}

    # 3) استبعد أعمدة ترقيم الصفوف والأعمدة الوصفية (المصدر/الملاحظات)
    rownum_cols = {c for c, cells in cols.items() if _is_rownum_col(cells)}
    rownum_cols |= ignore_cols

    # 4) رتّب الأعمدة حسب الأدوار (السعر ثم الكود ثم الاسم) — مع تفضيل تلميحات
    #    العناوين إن وُجدت. نبدأ بالسعر والكود لأن عناوين الأقسام قد تتسرّب إلى
    #    العمود الرقمي وتضخّم «نصّيته» فتخدع كاشف الاسم.
    name_scores = {c: _name_score(cells) for c, cells in cols.items()
                   if c not in rownum_cols}
    price_ratios = {c: _price_ratio(cells) for c, cells in cols.items()
                    if c not in rownum_cols}
    code_ratios = {c: _code_ratio(cells) for c, cells in cols.items()
                   if c not in rownum_cols}

    taken = set(rownum_cols)

    # عمود السعر (أعلى نسبة أرقام)
    price_col = hdr_map.get("price")
    if price_col is None or price_col in taken:
        cand = {c: r for c, r in price_ratios.items()
                if c not in taken and r >= 0.5}
        price_col = max(cand, key=cand.get) if cand else None
    if price_col is not None:
        taken.add(price_col)

    # عمود الكود (أعلى نسبة أكواد)
    code_col = hdr_map.get("code")
    if code_col is None or code_col in taken:
        cand = {c: r for c, r in code_ratios.items()
                if c not in taken and r >= 0.5}
        code_col = max(cand, key=cand.get) if cand else None
    if code_col is not None:
        taken.add(code_col)

    # عمود الاسم (أعلى «نصّية» بين الأعمدة المتبقية)
    name_col = hdr_map.get("name")
    if name_col is None or name_col in taken:
        cand = {c: s for c, s in name_scores.items() if c not in taken}
        name_col = max(cand, key=cand.get) if cand else None
    if name_col is not None:
        taken.add(name_col)

    # عمود الاسم البديل (ثاني أعلى عمود اسمي ذي محتوى نصّي معتبر)
    alt_col = None
    alt_cands = {c: name_scores.get(c, 0) for c in cols
                 if c not in taken and name_scores.get(c, 0) >= 3}
    if alt_cands:
        alt_col = max(alt_cands, key=alt_cands.get)

    if name_col is None:
        return []

    # 5) أنتج السجلات
    records = []
    have_price_col = price_col is not None
    for r in range(start, len(grid)):
        row = grid[r]
        name = row[name_col] if name_col < len(row) else ""
        alt = row[alt_col] if (alt_col is not None and alt_col < len(row)) else ""
        code = row[code_col] if (code_col is not None and code_col < len(row)) else ""
        price_raw = row[price_col] if (have_price_col and price_col < len(row)) else ""

        # لو الاسم فارغ والبديل موجود، بدّلهما
        if not name and alt:
            name, alt = alt, ""
        if not name:
            continue
        # السعر: رقمي إن أمكن، وإلا مرّر النص كما هو (مثل «35-650») للمراجع
        price = is_price(price_raw) or price_raw
        # تجاوز عناوين الأقسام: اسم موجود وخانة السعر **فارغة تمامًا**
        # وبلا كود وبلا اسم بديل (وعمود السعر معرّف)
        if have_price_col and not price_raw and not is_code(code) and not alt:
            # غالبًا عنوان قسم؛ تأكد أنه ليس خدمة مجانية مرقّمة بكود
            if not is_code(name):
                continue
        # لا تجعل الكود يتسرّب لخانة الاسم
        if is_code(name) and alt:
            name, alt = alt, name
        records.append({
            "name": name,
            "name_alt": alt if alt and alt != name else "",
            "code": code if is_code(code) else "",
            "price": price,
            "source": source,
        })
    return records


# ===========================================================================
# قُرّاء الصيغ
# ===========================================================================
# أوراق وصفية (ملخصات/مصادر) لا تحوي خدمات — تُتجاهل عند القراءة
_SKIP_SHEET = re.compile(r"ملخص|مصادر|المصادر|summary|source", re.IGNORECASE)


def _read_excel(path):
    sheets = pd.read_excel(path, sheet_name=None, header=None, dtype=str)
    out = []
    for sheet_name, df in sheets.items():
        if _SKIP_SHEET.search(str(sheet_name)):
            continue
        grid = df.fillna("").astype(str).values.tolist()
        out.extend(extract_records_from_grid(grid, sheet_name))
    return out


def _read_csv(path):
    # جرّب عدة فواصل تلقائيًا
    for sep in (None, ",", ";", "\t", "|"):
        try:
            df = pd.read_csv(path, header=None, dtype=str, sep=sep,
                             engine="python", encoding="utf-8-sig")
            if df.shape[1] >= 1:
                grid = df.fillna("").astype(str).values.tolist()
                return extract_records_from_grid(grid, os.path.basename(path))
        except Exception:
            continue
    return []


def _read_pdf(path):
    import pdfplumber

    records = []
    with pdfplumber.open(path) as pdf:
        for pi, page in enumerate(pdf.pages, 1):
            src = f"صفحة {pi}"
            tables = []
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            page_recs = []
            for t in tables:
                grid = [[fix_arabic(_clean(c)) for c in (row or [])] for row in t]
                page_recs.extend(extract_records_from_grid(grid, src))
            # احتياط نصّي: إن لم تُكتشف جداول مفيدة، حلّل الأسطر
            if not page_recs:
                page_recs = _parse_pdf_text(page, src)
            records.extend(page_recs)
    return records


_LINE_PRICE = re.compile(r"^(.*?)[\s.:]*?(\d[\d,،.\s]*\d|\d)\s*$")


def _parse_pdf_text(page, src):
    recs = []
    try:
        text = page.extract_text() or ""
    except Exception:
        return recs
    for line in text.splitlines():
        line = fix_arabic(_clean(line))
        if not line or not has_letter(line):
            continue
        m = _LINE_PRICE.match(line)
        if not m:
            continue
        name = m.group(1).strip(" .:-")
        price = is_price(m.group(2)) or ""
        if not name or not has_letter(name):
            continue
        recs.append({"name": name, "name_alt": "", "code": "",
                     "price": price, "source": src})
    return recs


def _read_pptx(path):
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(path)
    records = []
    for si, slide in enumerate(prs.slides, 1):
        src = f"شريحة {si}"
        had_table = False
        textboxes = []     # (top, left, text)
        for shape in slide.shapes:
            if shape.has_table:
                had_table = True
                tbl = shape.table
                grid = [[fix_arabic(_clean(cell.text)) for cell in row.cells]
                        for row in tbl.rows]
                records.extend(extract_records_from_grid(grid, src))
            elif shape.has_text_frame:
                txt = fix_arabic(_clean(shape.text_frame.text))
                if txt:
                    top = shape.top if shape.top is not None else 0
                    left = shape.left if shape.left is not None else 0
                    textboxes.append((int(top), int(left), txt))
        # مربّعات نصّية منفصلة: زاوج كل سعر بأقرب اسم موضعيًا
        if not had_table and textboxes:
            records.extend(_pair_textboxes(textboxes, src))
    return records


def _pair_textboxes(boxes, src):
    names = [(t, l, x) for (t, l, x) in boxes if has_letter(x)
             and is_price(x) is None]
    prices = [(t, l, x) for (t, l, x) in boxes if is_price(x) is not None]
    recs = []
    used = set()
    for (pt, pl, px) in prices:
        best, bestd = None, None
        for i, (nt, nl, nx) in enumerate(names):
            if i in used:
                continue
            d = (pt - nt) ** 2 + (pl - nl) ** 2
            if bestd is None or d < bestd:
                bestd, best = d, i
        if best is not None:
            used.add(best)
            recs.append({"name": names[best][2], "name_alt": "", "code": "",
                         "price": is_price(px) or "", "source": src})
    # أسماء بلا سعر مجاور تبقى كخدمات للمراجعة
    for i, (nt, nl, nx) in enumerate(names):
        if i not in used and len(nx) >= 3:
            recs.append({"name": nx, "name_alt": "", "code": "",
                         "price": "", "source": src})
    return recs


# ===========================================================================
# الواجهة العامة
# ===========================================================================
def load_records(path) -> pd.DataFrame:
    ext = os.path.splitext(path)[1].lower()
    if ext in (".xlsx", ".xls", ".xlsm"):
        records = _read_excel(path)
    elif ext == ".csv":
        records = _read_csv(path)
    elif ext == ".pdf":
        records = _read_pdf(path)
    elif ext in (".pptx", ".ppt"):
        records = _read_pptx(path)
    else:
        raise ValueError(f"صيغة غير مدعومة: {ext}")

    if not records:
        return pd.DataFrame(columns=OUT_COLS)
    df = pd.DataFrame(records)
    for c in OUT_COLS:
        if c not in df.columns:
            df[c] = ""
    df = df[OUT_COLS]
    # أزل التكرار التام والصفوف بلا اسم
    df = df[df["name"].astype(str).str.strip() != ""]
    df = df.drop_duplicates(subset=["name", "code", "price", "source"]).reset_index(drop=True)
    return df


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        sys.exit("الاستخدام: python ingest.py <ملف>")
    d = load_records(sys.argv[1])
    print(f"سجلات: {len(d)}")
    print(d.head(20).to_string())
